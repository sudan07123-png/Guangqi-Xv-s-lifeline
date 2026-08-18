# -*- coding: utf-8 -*-
"""
生成《徐光启 · 会通之志》作品介绍 PPT（比赛/评审用）
深色主题，配色与网站保持一致（墨绿 + 金 + 羊皮纸）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

ASSET = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_ppt_assets")

# ===== 配色（取自 styles/variables.css）=====
BG        = RGBColor(0x0e, 0x16, 0x11)   # 极深墨绿黑
BG_GREEN  = RGBColor(0x1b, 0x3d, 0x1c)   # --color-green
GREEN_MID = RGBColor(0x2d, 0x5a, 0x2e)   # --color-green-light
GOLD      = RGBColor(0xc9, 0xa9, 0x6e)   # --color-gold
GOLD_LT   = RGBColor(0xdf, 0xc4, 0x8a)   # --color-gold-light
PARCH     = RGBColor(0xf5, 0xf2, 0xeb)   # --color-parchment
MUTED     = RGBColor(0x9a, 0x96, 0x8e)   # --color-ink-muted
RED_MING  = RGBColor(0x8b, 0x1a, 0x1a)   # --color-red-ming
BLUE_WEST = RGBColor(0x4a, 0x6f, 0xa5)   # --color-blue-west
CARD      = RGBColor(0x16, 0x24, 0x1b)   # 卡片底色（略浅于背景）

FONT = "Microsoft YaHei"
FONT_SERIF = "SimSun"


def _set_run(run, text, size=16, color=PARCH, bold=False, font=FONT, italic=False):
    """设置 run 的文本与样式（含中文字体 east-asian 处理）"""
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn('a:rFonts'))
    if rFonts is None:
        rFonts = rPr.makeelement(qn('a:rFonts'), {})
        rPr.append(rFonts)
    rFonts.set(qn('a:ea'), font)
    rFonts.set(qn('a:cs'), font)
    rFonts.set(qn('a:latin'), font)


def add_text(slide, left, top, width, height, text, size=16, color=PARCH,
             bold=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        _set_run(p.add_run(), line, size=size, color=color, bold=bold, font=font)
    return box


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, shape=MSO_SHAPE.RECTANGLE, line=False):
    sp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = GOLD
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_pic(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return None
    kw = {}
    if width:
        kw["width"] = Inches(width)
    if height:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)


def header(slide, title, subtitle=None, idx=None):
    """统一页眉：金色细线 + 标题"""
    add_rect(slide, 0.9, 0.62, 0.42, 0.045, GOLD)
    add_text(slide, 0.9, 0.75, 11.5, 0.7, title, size=27, color=PARCH, bold=True)
    if subtitle:
        add_text(slide, 0.92, 1.38, 11.0, 0.4, subtitle, size=12.5, color=MUTED)
    if idx:
        add_text(slide, 12.2, 0.72, 0.7, 0.5, str(idx).zfill(2), size=20,
                 color=GOLD, bold=True, align=PP_ALIGN.RIGHT)


def footer(slide, n):
    add_text(slide, 0.9, 7.02, 6.0, 0.3, "徐光启 · 会通之志  作品介绍", size=9, color=MUTED)
    add_text(slide, 12.0, 7.02, 0.9, 0.3, f"{n:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def bullet_box(slide, left, top, width, height, items):
    """items: list of (标题, 描述) 或 (标题, None)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for title, desc in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.line_spacing = 1.1
        _set_run(p.add_run(), "▍", size=14, color=GOLD, bold=True)
        _set_run(p.add_run(), title + "  ", size=15, color=PARCH, bold=True)
        if desc:
            _set_run(p.add_run(), desc, size=12.5, color=MUTED)
    return box


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============================================================
# 1. 封面
# ============================================================
s = new_slide()
add_bg(s, BG_GREEN)
# 上下金色细框
add_rect(s, 0, 0, SW, 0.14, GOLD)
add_rect(s, 0, SH - 0.14, SW, 0.14, GOLD)
add_rect(s, 0.9, 2.05, 0.06, 1.5, GOLD)
add_text(s, 1.15, 2.05, 10.5, 0.4, "交互式历史叙事 · 微站作品", size=15, color=GOLD_LT)
add_text(s, 1.12, 2.45, 11.0, 1.6, "徐光启 · 会通之志", size=60, color=PARCH, bold=True, font=FONT_SERIF)
add_text(s, 1.18, 3.95, 10.5, 0.5, "一个关于知识、信仰与抉择的大明故事", size=20, color=PARCH)
# 分隔
add_rect(s, 1.18, 4.62, 3.2, 0.03, GOLD)
add_text(s, 1.18, 4.85, 11.0, 0.4, "徐光启（1562–1633） · 中西会通第一人 · 译《几何原本》 · 修《崇祯历书》",
         size=13, color=MUTED)
add_text(s, 1.18, 6.55, 11.0, 0.4, "—— 参赛作品 · 交互式历史叙事 ——", size=12, color=GOLD_LT, align=PP_ALIGN.CENTER)

# ============================================================
# 2. 作品简介
# ============================================================
s = new_slide()
add_bg(s)
header(s, "作品简介", "以徐光启一生为蓝本的交互式历史叙事游戏", 2)
add_rect(s, 0.9, 2.05, 6.7, 4.4, CARD)
add_text(s, 1.25, 2.4, 6.0, 0.4, "定位", size=13, color=GOLD, bold=True)
add_text(s, 1.25, 2.85, 6.0, 2.5,
         "玩家化身明代科学家徐光启，从上海法华汇的寒门少年，一路走到崇祯年间主持修历。"
         "在「科举入仕」与「钻研西学」之间反复抉择，亲历那段中西文明碰撞的大时代。",
         size=15, color=PARCH, line_spacing=1.4)
add_text(s, 1.25, 5.35, 6.0, 0.9,
         "一部可以「玩」的历史课——让用户在体验中理解先贤的孤独与坚守。",
         size=12.5, color=MUTED, line_spacing=1.3)
# 右侧数据卡片
data = [("3 章", "十七个叙事场景"), ("5 节点", "关键抉择节点"), ("12 条", "支线走向"), ("6 种", "结局（含隐藏）")]
for i, (big, small) in enumerate(data):
    x = 7.95 + (i % 2) * 2.35
    y = 2.05 + (i // 2) * 2.35
    add_rect(s, x, y, 2.15, 2.05, BG_GREEN)
    add_rect(s, x, y, 2.15, 0.06, GOLD)
    add_text(s, x + 0.2, y + 0.35, 1.8, 0.8, big, size=30, color=GOLD_LT, bold=True)
    add_text(s, x + 0.2, y + 1.2, 1.8, 0.4, small, size=11.5, color=PARCH)
footer(s, 2)

# ============================================================
# 3. 故事与历史背景
# ============================================================
s = new_slide()
add_bg(s)
header(s, "故事背景", "为什么选择徐光启？", 3)
items = [
    ("「中西会通第一人」", "中国最早系统引入西方科学（天文、数学、地理）的学者之一，与利玛窦合译《几何原本》，奠定了现代科学汉译词汇。"),
    ("从寒门到阁臣", "出身上海农家的书生，历经科考坎坷，最终官至礼部尚书、内阁大学士，一生充满张力与抉择。"),
    ("会通以超胜", "他提出的「会通中西、会通超胜」思想，至今仍是文明互鉴的深刻命题。"),
]
bullet_box(s, 1.0, 2.1, 7.0, 4.2, items)
add_rect(s, 8.3, 2.1, 4.0, 4.3, CARD)
add_pic(s, os.path.join(ASSET, "title.png"), 8.5, 2.3, width=3.6)
add_text(s, 8.5, 5.6, 3.6, 0.6, "标题页界面", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 3)

# ============================================================
# 4. 核心玩法
# ============================================================
s = new_slide()
add_bg(s)
header(s, "核心玩法", "交互式历史叙事：在阅读中「参与」历史", 4)
items = [
    ("章节叙事 · 打字机演出", "三章十七个场景，以旁白、对白、内心独白逐幕展开，营造沉浸的阅读节奏。"),
    ("关键抉择 · 分支推进", "五个抉择节点贯穿全程，每一次选择都会写入状态旗标（flags），影响后续剧情走向。"),
    ("多结局 · 可重玩", "不同选择组合通向六种结局，鼓励玩家反复探索，理解「选择塑造人生」。"),
]
bullet_box(s, 1.0, 2.1, 6.7, 4.3, items)
add_rect(s, 8.0, 2.1, 4.3, 4.3, CARD)
add_pic(s, os.path.join(ASSET, "narrative_choice.png"), 8.15, 2.25, width=4.0)
add_text(s, 8.15, 5.65, 4.0, 0.6, "抉择界面（含选项按钮）", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 4)

# ============================================================
# 5. 章节结构
# ============================================================
s = new_slide()
add_bg(s)
header(s, "章节结构", "三段人生，一条主线", 5)
chapters = [
    ("第一章 · 寒门少年", "1562 — 1600", "出生法华汇 → 辗转两广 → 韶州偶遇世界地图 → 北上中举", "从泥土与书本之间走出，被一张世界地图点燃。"),
    ("第二章 · 译书与功名", "1600 — 1613", "南京初遇利玛窦 → 进士及第 → 翰林院译书 → 师丧", "功名与西学并进，译《几何原本》于会通之际。"),
    ("第三章 · 归去来兮", "1613 — 1633", "党争 → 天津屯田 → 萨尔浒 → 历法改革 → 最终抉择", "在朝堂与田野之间，走向修历的最后一程。"),
]
for i, (t, sub, path, desc) in enumerate(chapters):
    y = 2.05 + i * 1.62
    add_rect(s, 0.9, y, 11.5, 1.4, CARD if i % 2 == 0 else BG_GREEN)
    add_rect(s, 0.9, y, 0.08, 1.4, GOLD)
    add_text(s, 1.25, y + 0.18, 3.0, 0.6, t, size=18, color=PARCH, bold=True)
    add_text(s, 4.45, y + 0.24, 1.6, 0.5, sub, size=13, color=GOLD, bold=True)
    add_text(s, 1.25, y + 0.75, 10.9, 0.5, path + "　｜　" + desc, size=12, color=MUTED)
footer(s, 5)

# ============================================================
# 6. 分支与结局系统
# ============================================================
s = new_slide()
add_bg(s)
header(s, "分支与结局系统", "5 节点 · 12 支线 · 6 结局", 6)
add_text(s, 0.9, 2.05, 11.5, 0.9,
         "节点间通过「状态旗标（flags）」连锁：前一节点的选择，会开启或关闭后续节点的选项，"
         "最终汇聚为不同的结局。", size=14, color=PARCH, line_spacing=1.3)
# 节点链条
nodes = ["① 韶州地图", "② 南京利玛窦", "③ 翰林院译书", "④ 天津屯田", "⑤ 修历"]
for i, n in enumerate(nodes):
    x = 0.9 + i * 2.38
    add_rect(s, x, 3.1, 2.05, 0.62, BG_GREEN)
    add_text(s, x, 3.22, 2.05, 0.4, n, size=12.5, color=GOLD_LT, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, x + 2.05, 3.22, 0.33, 0.4, "→", size=15, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.95, 11.5, 0.4, "隐藏结局「大明科学院」：①A → ②A → ③B → ④A → ⑤A 全对通关方可解锁",
         size=12.5, color=GOLD_LT)
# 结局卡
endings = [
    ("启明孤星", "真结局", GOLD), ("大明科学院", "隐藏结局", GOLD_LT),
    ("良史未竟", "普通结局", PARCH), ("归隐田园", "普通结局", PARCH),
    ("鞠躬尽瘁", "遗憾结局", MUTED), ("庸臣之叹", "失败结局", MUTED),
]
for i, (t, sub, c) in enumerate(endings):
    x = 0.9 + (i % 3) * 3.95
    y = 4.7 + (i // 3) * 1.05
    add_rect(s, x, y, 3.65, 0.85, CARD, line=True)
    add_text(s, x + 0.22, y + 0.12, 2.4, 0.4, t, size=14.5, color=c, bold=True)
    add_text(s, x + 0.22, y + 0.5, 3.2, 0.3, sub, size=10.5, color=MUTED)
footer(s, 6)

# ============================================================
# 7. 创新点一 · 史料考据引用系统
# ============================================================
s = new_slide()
add_bg(s)
header(s, "创新点一 · 史料考据引用系统", "让叙事成为「可信的历史」", 7)
add_text(s, 0.9, 2.0, 11.5, 1.1,
         "叙事文本中的关键史实（人名、地名、事件）以下划线高亮，鼠标悬停即弹出考据出处，"
         "溯源至《明史》《明实录》等原始史料。", size=14, color=PARCH, line_spacing=1.35)
add_rect(s, 0.9, 3.2, 11.5, 2.6, CARD)
add_text(s, 1.25, 3.45, 10.8, 0.5, "示例 · 第一章开场旁白", size=12, color=GOLD, bold=True)
add_text(s, 1.25, 3.95, 10.8, 1.5,
         "「嘉靖四十一年（1562），上海县法华汇，一户务农兼经商的人家，添了一个男孩。」\n"
         "—— 其中「上海县法华汇」带下划线，悬停显示引文出处与原文。",
         size=14.5, color=PARCH, line_spacing=1.4)
add_text(s, 0.9, 6.0, 11.5, 0.6, "价值：将游戏从「戏说」提升为「寓教于乐」的历史学习工具。",
         size=13, color=GOLD_LT)
footer(s, 7)

# ============================================================
# 8. 创新点二 · 时空朋友圈
# ============================================================
s = new_slide()
add_bg(s)
header(s, "创新点二 · 时空朋友圈", "以现代社交形式，重现明代人物的交游网络", 8)
add_text(s, 0.9, 2.0, 6.3, 1.0,
         "将徐光启与利玛窦、家人、同僚的关系，映射为「朋友圈」时间轴，"
         "用当代年轻人熟悉的界面语言，降低理解历史人物的门槛。",
         size=14, color=PARCH, line_spacing=1.35)
items = [
    ("关系可视化", "徐父、徐母、利玛窦、徐骥、崇祯、保守派朝臣——关系亲疏一目了然。"),
    ("时间轴联动", "随着章节推进，朋友圈内容与剧情同步更新。"),
]
bullet_box(s, 0.9, 3.2, 6.3, 3.0, items)
add_rect(s, 7.5, 2.0, 4.9, 4.5, CARD)
add_pic(s, os.path.join(ASSET, "friends.png"), 7.65, 2.15, width=4.6)
add_text(s, 7.65, 5.75, 4.6, 0.6, "时空朋友圈界面", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 8)

# ============================================================
# 9. 创新点三 · 地图枢纽
# ============================================================
s = new_slide()
add_bg(s)
header(s, "创新点三 · 时空地图枢纽", "在「时空漫游」中串联历史足迹", 9)
add_text(s, 0.9, 2.0, 6.3, 1.0,
         "以地图 + 生平大事记的形式，将徐光启的人生足迹（上海、韶州、南京、北京、天津）"
         "与章节选择整合为可视化枢纽。",
         size=14, color=PARCH, line_spacing=1.35)
items = [
    ("足迹可视化", "地图标记与时间轴对应，串联人物行迹。"),
    ("章节导航", "底部章节卡片作为进入各章的入口，兼顾叙事与总览。"),
]
bullet_box(s, 0.9, 3.2, 6.3, 3.0, items)
add_rect(s, 7.5, 2.0, 4.9, 4.5, CARD)
add_pic(s, os.path.join(ASSET, "map.png"), 7.65, 2.15, width=4.6)
add_text(s, 7.65, 5.75, 4.6, 0.6, "地图枢纽界面", size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 9)

# ============================================================
# 10. 技术实现
# ============================================================
s = new_slide()
add_bg(s)
header(s, "技术实现", "纯前端 · 零依赖 · 数据驱动", 10)
techs = [
    ("单页应用（SPA）", "基于 hash 路由（#/narrative/…），无需刷新即可在标题、叙事、朋友圈、地图、结局五页间切换。"),
    ("IIFE 模块化架构", "事件总线（EventBus）+ 状态管理（StateManager）+ 路由（Router）分层解耦，核心/引擎/UI 三组模块清晰。"),
    ("JSON 数据驱动", "章节、人物、支线、结局、史料引用全部由 JSON 配置，剧情与内容可脱离代码维护。"),
    ("相对路径 · 即开即用", "所有资源均用相对路径，可直接部署到 GitHub Pages 等任意静态托管，双击 start.bat 本地即启。"),
]
for i, (t, d) in enumerate(techs):
    y = 2.0 + i * 1.28
    add_rect(s, 0.9, y, 11.5, 1.1, CARD if i % 2 == 0 else BG_GREEN)
    add_text(s, 1.25, y + 0.16, 3.4, 0.5, t, size=15, color=GOLD_LT, bold=True)
    add_text(s, 1.25, y + 0.6, 10.9, 0.5, d, size=12, color=MUTED)
footer(s, 10)

# ============================================================
# 11. 界面展示
# ============================================================
s = new_slide()
add_bg(s)
header(s, "界面展示", "五页 · 五种体验", 11)
shots = [
    ("title.png", "标题页", 0.9, 2.1),
    ("narrative_choice.png", "叙事·抉择", 5.25, 2.1),
    ("map.png", "地图枢纽", 0.9, 4.45),
    ("friends.png", "时空朋友圈", 5.25, 4.45),
    ("ending.png", "结局页", 9.6, 2.1),
]
for name, label, x, y in shots:
    add_rect(s, x - 0.08, y - 0.08, 4.05, 2.42, CARD)
    add_pic(s, os.path.join(ASSET, name), x, y, width=3.9)
    add_text(s, x, y + 2.1, 3.9, 0.3, label, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 11)

# ============================================================
# 12. 访问与部署
# ============================================================
s = new_slide()
add_bg(s)
header(s, "访问与部署", "一个网址即可体验", 12)
add_rect(s, 0.9, 2.0, 11.5, 1.3, BG_GREEN)
add_text(s, 1.25, 2.2, 3.0, 0.5, "公网访问", size=13, color=GOLD, bold=True)
add_text(s, 1.25, 2.65, 10.8, 0.5, "https://<用户名>.github.io/<仓库名>/", size=16, color=PARCH, bold=True)
steps = [
    ("本地预览", "双击根目录 start.bat（或运行 python -m http.server 8080），浏览器访问 http://localhost:8080"),
    ("部署到 GitHub Pages", "将项目推送到 GitHub 仓库 → Settings → Pages → 选择分支与根目录 → 保存，等待生成公网地址。"),
    ("无需后端", "纯静态站点，无服务器、无数据库，部署零成本、长期可访问。"),
]
for i, (t, d) in enumerate(steps):
    y = 3.65 + i * 1.05
    add_text(s, 1.0, y, 1.1, 0.4, f"0{i+1}", size=18, color=GOLD, bold=True)
    add_text(s, 2.1, y + 0.02, 2.2, 0.4, t, size=14, color=PARCH, bold=True)
    add_text(s, 4.4, y + 0.02, 8.0, 0.5, d, size=12, color=MUTED)
footer(s, 12)

# ============================================================
# 13. 总结与展望
# ============================================================
s = new_slide()
add_bg(s, BG_GREEN)
add_rect(s, 0, 0, SW, 0.14, GOLD)
add_rect(s, 0, SH - 0.14, SW, 0.14, GOLD)
add_text(s, 0.9, 1.7, 11.5, 0.6, "总结", size=30, color=PARCH, bold=True)
add_text(s, 0.9, 2.5, 11.5, 2.0,
         "《徐光启 · 会通之志》不是一部传统的历史科普页面，而是一场「可体验的抉择」。"
         "它以交互叙事、史料考据、分支结局，让四百年前那位孤独的先驱重新「活」在当下——"
         "在每一次点击与选择中，传递「会通中西、会通超胜」的永恒价值。",
         size=16, color=PARCH, line_spacing=1.5)
add_rect(s, 0.9, 4.6, 3.2, 0.03, GOLD)
add_text(s, 0.9, 4.85, 11.5, 0.5, "展望", size=16, color=GOLD_LT, bold=True)
add_text(s, 0.9, 5.4, 11.5, 1.0,
         "后续 Phase 将持续完善：集成真实地图（Leaflet）、补齐人物立绘与背景美术、扩充支线剧情与音效，"
         "让作品更加完整与沉浸。", size=13, color=PARCH, line_spacing=1.4)
add_text(s, 0.9, 6.7, 11.5, 0.4, "谢谢评审老师 · 欢迎体验", size=13, color=GOLD_LT, align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "徐光启_会通之志_作品介绍.pptx")
prs.save(out)
print("已生成:", out)
print("总页数:", len(prs.slides._sldIdLst))
